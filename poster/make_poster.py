"""Generate the CIKM 2026 poster for PMC (csp4014) as an editable PPTX.

v3 — dense and self-contained, in the style of the CIKM / ICLR / ECCV posters
the authors referenced: an explicit Background, Contributions, Setup,
Conclusion and full explanatory prose, so a visitor can read the whole
argument without the presenter standing there.

Visual language follows the lab's own deck: white ground, KU-crimson serif
headings (#8B0029, sampled), hairline rules, IE-LAB blue (#0D4DA3) as the
single accent meaning PMC, gray for baselines. The official CIKM '26 mark
sits in the footer band with the conference details.

Layout is a flow model: each column owns a top-down cursor and every block
reports the height it consumed, so a block can never silently overrun its
neighbour. build() prints per-column fill against the available height.

A0 portrait (33.11 x 46.81 in); rescale once CIKM publishes the poster spec.

Usage:  python make_poster.py
"""

import re
from pathlib import Path

from PIL import ImageFont

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

HERE = Path(__file__).resolve().parent

CRIMSON = RGBColor(0x8B, 0x00, 0x29)
BLUE = RGBColor(0x0D, 0x4D, 0xA3)
INK = RGBColor(0x1F, 0x1F, 0x1F)
GRAY = RGBColor(0x6F, 0x6F, 0x6F)
HAIR = RGBColor(0xC8, 0xC8, 0xC8)
CARD = RGBColor(0xF6, 0xF5, 0xF3)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SERIF = "Times New Roman"

PAGE_W, PAGE_H = Inches(33.11), Inches(46.81)
MARGIN = Inches(0.85)
HEADER_H = Inches(5.25)
FOOTER_H = Inches(2.05)

ASPECT = {                        # native pixel aspect of each embedded asset
    "asset_fig1.png": 2400 / 1193,
    "asset_fig2ov.png": 1537 / 1024,
    "asset_fig3b.png": 1188 / 929,
    "asset_fig3c.png": 1206 / 929,
    "asset_gapgain.png": 989 / 657,
    "asset_fig3.png": 3616 / 933,
    "asset_ku.jpg": 157 / 49,
    "asset_ielab.png": 177 / 71,
    "asset_cikm.png": 250 / 193,
}

# Line breaking is measured with the real Times metrics (PIL) rather than
# estimated, so the height a block reserves matches what PowerPoint sets and
# no block can overlap its neighbour.
LINE_FACTOR = 1.20      # PowerPoint sets leading at 1.2x the font size
_TIMES = {False: "/System/Library/Fonts/Supplemental/Times New Roman.ttf",
          True: "/System/Library/Fonts/Supplemental/Times New Roman Bold.ttf"}
_MPX = 64               # metric render size; widths scale linearly from it
_FONTS = {}


def solid(shape, color, line_color=None, line_w=None):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
        shape.line.width = line_w or Pt(1.0)
    shape.shadow.inherit = False


def _width_pt(text, size, bold):
    if bold not in _FONTS:
        _FONTS[bold] = ImageFont.truetype(_TIMES[bold], _MPX)
    return _FONTS[bold].getlength(text) * size / _MPX


def wrap_lines(runs, width_pt):
    """Greedy word-wrap over mixed-format runs; returns each line's max size."""
    lines, cur_w, cur_max = [], 0.0, 0.0
    for text, size, bold in runs:
        for tok in re.findall(r"\s+|\S+", text):
            w = _width_pt(tok, size, bold)
            if tok.isspace():
                if cur_w:
                    cur_w += w
                continue
            if cur_w and cur_w + w > width_pt:
                lines.append(cur_max)
                cur_w, cur_max = w, size
            else:
                cur_w += w
                cur_max = max(cur_max, size)
    if cur_max:
        lines.append(cur_max)
    return lines or [runs[0][1]]


def add_runs(p, text, size, bold, color, italic=False):
    """Emit runs, rendering `_{...}` as a true subscript."""
    parts = re.split(r"_\{([^}]*)\}", text)
    for i, part in enumerate(parts):
        if not part:
            continue
        r = p.add_run()
        r.text = part
        r.font.size = Pt(size * 0.82) if i % 2 else Pt(size)
        r.font.bold, r.font.italic, r.font.name = bold, italic, SERIF
        r.font.color.rgb = color
        if i % 2:
            r.font._rPr.set("baseline", "-25000")


def plain(text):
    return re.sub(r"_\{([^}]*)\}", r"\1", text)


def _runs_of(para, size, bold, color, italic):
    return para if isinstance(para[0], tuple) else [(para, size, bold, color, italic)]


def cell(slide, runs, x, y, w, h, align=PP_ALIGN.LEFT, wrap=False):
    """One line of mixed runs, vertically centred — used by tables and bars."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    for t, sz, bd, cl in runs:
        add_runs(p, t, sz, bd, cl)
    return tb


class Col:
    """A poster column with a top-down cursor."""

    def __init__(self, slide, x, w, top, bottom):
        self.s, self.x, self.w, self.y, self.bottom = slide, x, w, top, bottom

    def gap(self, inches):
        self.y += Inches(inches)

    def rule(self, weight=Pt(1.0), color=HAIR, pad=0.10):
        self.y += Inches(pad)
        solid(self.s.shapes.add_shape(MSO_SHAPE.RECTANGLE, self.x, self.y, self.w, weight),
              color)
        self.y += weight + Inches(pad)

    def text(self, paras, size=21, color=INK, bold=False, italic=False,
             align=PP_ALIGN.LEFT, spacing=1.22, after=5, x=None, w=None):
        x = self.x if x is None else x
        w = self.w if w is None else w
        width_pt = w / 12700.0
        total_pt = 0.0
        for p in paras:
            runs = [(plain(r[0]), r[1] if len(r) > 1 else size,
                     r[2] if len(r) > 2 else bold)
                    for r in _runs_of(p, size, bold, color, italic)]
            total_pt += sum(ls * spacing * LINE_FACTOR
                            for ls in wrap_lines(runs, width_pt)) + after
        h = Inches(total_pt / 72.0)
        tb = self.s.shapes.add_textbox(x, self.y, w, h)
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        for i, para in enumerate(paras):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align
            p.line_spacing = spacing
            p.space_after = Pt(after)
            for run in _runs_of(para, size, bold, color, italic):
                add_runs(p, run[0],
                         run[1] if len(run) > 1 else size,
                         run[2] if len(run) > 2 else bold,
                         run[3] if len(run) > 3 else color,
                         run[4] if len(run) > 4 else italic)
        self.y += h
        return h

    def heading(self, num, title):
        self.gap(0.08)
        self.text([[(f"{num}.  {title}", 35, True, CRIMSON)]], spacing=1.05, after=0)
        self.rule(Pt(2.5), CRIMSON, pad=0.05)

    def bullets(self, items, size=21, marker=CRIMSON, gap_in=0.09):
        for it in items:
            solid(self.s.shapes.add_shape(MSO_SHAPE.RECTANGLE, self.x + Inches(0.03),
                                          self.y + Inches(0.15), Inches(0.18), Inches(0.18)),
                  marker)
            self.text([it], size=size, x=self.x + Inches(0.44),
                      w=self.w - Inches(0.44), spacing=1.20, after=0)
            self.gap(gap_in)

    def figure(self, name, frac=1.0):
        w = Emu(int(self.w * frac))
        h = Emu(int(w / ASPECT[name]))
        self.s.shapes.add_picture(str(HERE / name),
                                  self.x + Emu(int((self.w - w) / 2)), self.y, w, h)
        self.y += h

    def _boxed(self, paras, size, fill, pad, color, bold, line):
        # The box is created before the text so it sits behind it in z-order;
        # its height is corrected once the text has advanced the cursor.
        start = self.y
        box = self.s.shapes.add_shape(MSO_SHAPE.RECTANGLE, self.x, start, self.w,
                                      Inches(1))
        solid(box, fill, line, Pt(1.0) if line else None)
        self.gap(pad)
        self.text(paras, size=size, color=color, bold=bold, x=self.x + Inches(pad),
                  w=self.w - Inches(2 * pad), spacing=1.28, after=4)
        self.gap(pad)
        box.height = self.y - start

    def card(self, paras, size=21, pad=0.28):
        self._boxed(paras, size, CARD, pad, INK, False, HAIR)

    def band(self, paras, size=24, pad=0.26):
        self._boxed(paras, size, BLUE, pad, WHITE, True, None)

    def hbar(self, frac, color, label, value, h_in=0.68, val_w=1.70):
        h = Inches(h_in)
        lane_w = self.w - Inches(val_w)
        solid(self.s.shapes.add_shape(MSO_SHAPE.RECTANGLE, self.x, self.y, lane_w, h),
              WHITE, HAIR, Pt(1.0))
        solid(self.s.shapes.add_shape(MSO_SHAPE.RECTANGLE, self.x, self.y,
                                      Emu(int(lane_w * frac)), h), color)
        cell(self.s, [(label, 24, False, WHITE if frac > 0.45 else INK)],
             self.x + Inches(0.18), self.y, lane_w - Inches(0.3), h)
        cell(self.s, [(value, 24, True, INK)],
             self.x + lane_w + Inches(0.14), self.y, Inches(val_w - 0.14), h)
        self.y += h + Inches(0.10)


ABSTRACT = (
    "Approximate nearest neighbor search is a core operator in large-scale retrieval, where"
    " vector quantization is widely used to reduce memory cost. Binary quantization methods"
    " encode each dimension as a single bit, but they assume queries and database vectors s"
    "hare a common centroid. This assumption fails in cross-modal retrieval, where each mod"
    "ality clusters around a different centroid, producing a modality gap that corrupts the"
    " centroid's dual role as inverted-file routing pivot and sign-bit decision boundary. W"
    "e propose Per-Modality Centroid Correction (PMC), which shifts database vectors toward"
    " the query centroid at build time, rewriting routing and code boundaries at the source"
    " rather than merely adjusting queries at serving time. A selective variant confirms th"
    "at concentrated gap energy in a few dimensions drives most recall loss. PMC adds no in"
    "dex memory or serving overhead; at α=1, the correction is fully absorbed into the inde"
    "x with zero query-time cost. Experiments on five cross-modal benchmarks, three encoder"
    "s, and four binary-quantized index types show gains over uncorrected and query-shifted"
    " baselines that scale with modality-gap magnitude, up to a 400M-vector deployment.")


def build():
    prs = Presentation()
    prs.slide_width, prs.slide_height = PAGE_W, PAGE_H
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    solid(slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, PAGE_W, PAGE_H), WHITE)
    M, CW = MARGIN, PAGE_W - 2 * MARGIN

    # ============================= header =============================
    KU_H = IE_H = Inches(1.15)
    KU_W = Emu(int(KU_H * ASPECT["asset_ku.jpg"]))
    IE_W = Emu(int(IE_H * ASPECT["asset_ielab.png"]))
    R = PAGE_W - M
    slide.shapes.add_picture(str(HERE / "asset_ielab.png"), R - IE_W, Inches(0.65), IE_W, IE_H)
    slide.shapes.add_picture(str(HERE / "asset_ku.jpg"),
                             R - IE_W - Inches(0.5) - KU_W, Inches(0.65), KU_W, KU_H)
    QR = Inches(2.05)
    slide.shapes.add_picture(str(HERE / "asset_qr.png"), R - QR, Inches(2.20), QR, QR)
    cell(slide, [("Paper + code:  github.com/sehoon787/PMC", 17, False, GRAY)],
         R - Inches(6.5), Inches(4.35), Inches(6.5), Inches(0.45), PP_ALIGN.RIGHT)

    head = Col(slide, M, CW - Inches(9.0), Inches(0.72), Inches(5.1))
    head.text([[("PMC: Build-Time Per-Modality Centroid Correction", 55, True, CRIMSON)],
               [("for Cross-Modal Binary-Quantized Retrieval", 55, True, CRIMSON)]],
              spacing=1.10, after=0)
    head.gap(0.26)
    head.text([[("Se Hoon Kim,   Jun Hyung Lee,   Soonyoung Jung*", 30, False, INK)]],
              after=0)
    head.gap(0.08)
    head.text([[("Department of Computer Science and Engineering, Korea University  ·  "
                 "Intelligence Engineering Lab", 22, False, GRAY)]], after=0)
    head.gap(0.06)
    head.text([[("{sehoon787, junicus, jsy}@korea.ac.kr", 20, False, GRAY, True)]], after=0)
    solid(slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, M, HEADER_H, CW, Pt(3.0)), CRIMSON)

    # ================= two-column split =================================
    TOP = HEADER_H + Inches(0.42)
    BOT = PAGE_H - FOOTER_H - Inches(0.30)
    GUT = Inches(0.70)
    W2 = Emu(int((CW - GUT) / 2))
    L = Col(slide, M, W2, TOP, BOT)
    R = Col(slide, M + W2 + GUT, W2, TOP, BOT)
    solid(slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 R.x - Emu(int(GUT / 2)), TOP, Pt(1.0), BOT - TOP),
          HAIR)

    def equation(col, runs, size=28, note=None):
        col.gap(0.16)
        col.text([runs], size=size, align=PP_ALIGN.CENTER, after=0)
        if note:
            col.text([[(note, 20, False, GRAY, True)]], size=20,
                     align=PP_ALIGN.CENTER, after=0)
        col.gap(0.18)

    # ---------------------------- LEFT ---------------------------------
    L.card([[("Abstract.  ", 22, True, CRIMSON), (ABSTRACT, 22, False, INK)]],
           size=22, pad=0.30)
    L.gap(0.22)

    L.heading("1", "The problem: the modality gap")
    L.text([[("A frozen multimodal encoder φ maps both modalities into d-dimensional space: the "
              "database holds modality-b embeddings, queries come from modality a. "
              "Their means differ — the ", 22, False, INK),
             ("modality gap", 22, True, CRIMSON),
             ("  g = μ_{q} − μ_{x}  is stable and large, ‖g‖ = .61–.82 on our "
              "benchmarks.", 22, False, INK)]], size=22)
    L.gap(0.14)
    L.figure("asset_fig1.png", 0.93)
    L.gap(0.08)
    L.text([[("t-SNE of ImageBind embeddings: each modality forms its own cluster; "
              "after PMC the paired points overlap.", 20, False, GRAY, True)]],
           size=20, spacing=1.15, after=0)
    L.gap(0.20)
    L.text([[("BQ encodes each dimension as  sign(x_{i} − c_{i})  relative to a "
              "centroid c (RaBitQ, Lucene BBQ). With one bit the decision boundary "
              "passes exactly through c with zero error margin — and the same "
              "centroid routes queries to IVF lists, so one offset corrupts codes "
              "and routing at once. Let δ_{i} = x_{i} − c_{i} be the residual; "
              "shifting the centroid by αg flips dimension i's sign iff",
              22, False, INK)]], size=22)
    equation(L, [("δ_{i} (δ_{i} − α g_{i})  <  0 ,      i.e.,      "
                  "0  <  δ_{i} / (α g_{i})  <  1", 27, True, INK)],
             note="δ_{i}: residual x_{i} − c_{i}   ·   g_{i}: gap component   ·   "
                  "α: shift strength")
    L.text([[("Under smooth symmetric residual densities, a first-order expansion "
              "gives the expected number of flipped bits,", 22, False, INK)]],
           size=22)
    equation(L, [("E[F]  ≈  α Σ_{i} p_{i}(0) |g_{i}|", 30, True, INK)],
             note="E[F]: expected bit flips   ·   p_{i}(0): residual density at "
                  "the code boundary")
    L.text([[("a density-weighted ℓ₁ norm of the gap: flip risk grows with |g_{i}|. "
              "On CLIP-L the top 10% of dimensions carry ≈90% of ‖g‖², so the "
              "corruption is structured rather than random, biasing Hamming "
              "distances systematically.", 22, False, INK)]], size=22)
    L.gap(0.14)
    L.card([[("Oracle test — MSCOCO, CLIP-L, RaBitQ.   ", 22, True, INK),
             ("Same-modality queries (‖g‖ = 0) reach R@100 = 0.71; cross-modal "
              "queries (‖g‖ = .82) fall to ", 22, False, INK),
             ("0.54", 22, True, CRIMSON),
             (". The loss is the gap, not the quantizer.", 22, False, INK)]],
           size=22, pad=0.26)

    L.heading("2", "Method: PMC correction")
    L.text([[("PMC applies a one-time centroid alignment before index "
              "construction. From a small paired sample we compute g and shift "
              "database vectors toward the query centroid:", 22, False, INK)]],
           size=22)
    equation(L, [("x′ = (x + α g) / ‖x + α g‖ ,      "
                  "q′ = (q − (1−α) g) / ‖q − (1−α) g‖", 27, True, INK)],
             note="x: database vector   ·   q: query   ·   g: modality gap   ·"
                  "   α ∈ [0,1]: shift strength  (α = 1  →  q′ = q)")
    L.text([[("α interpolates between full query-side correction (α = 0, mean "
              "shift) and full database-side alignment (α = 1). The BQ index is "
              "built on {x′}; at α = 1 the query-time transform is identity, so "
              "serving incurs no extra cost. Per-vector renormalization preserves "
              "the unit-norm structure, and the α-sweep confirms α = 1 is best or "
              "near-best in every setting.", 22, False, INK)]], size=22)
    L.gap(0.16)

    # Algorithm 1 — exactly as in the paper
    alg_top = L.y
    solid(slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, L.x, L.y, L.w, Pt(2.2)), INK)
    L.gap(0.10)
    L.text([[("Algorithm 1  PMC Index Construction", 23, True, INK)]], after=0)
    L.gap(0.08)
    solid(slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, L.x, L.y, L.w, Pt(1.0)), INK)
    L.gap(0.10)
    L.text([[("Require:", 21, True, INK),
             ("  database X (modality b), query sample Q (modality a), "
              "interpolation α", 21, False, INK)],
            [("Ensure:", 21, True, INK),
             ("  IVF-based quantized index I and gap g", 21, False, INK)],
            [("1:  μ_{q} ← mean(Q);   μ_{x} ← mean(X)", 21, False, INK)],
            [("2:  g ← μ_{q} − μ_{x}", 21, False, INK)],
            [("3:  for each x ∈ X do", 21, False, INK)],
            [("4:      x′ ← normalize(x + α g)", 21, False, INK)],
            [("5:  end for", 21, False, INK)],
            [("6:  I ← BuildIndex({x′})        — RaBitQ · BBQ · BinaryFlat",
              21, False, INK)],
            [("7:  return I, g", 21, False, INK)]],
           size=21, spacing=1.30, after=3)
    L.gap(0.08)
    solid(slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, L.x, L.y, L.w, Pt(2.2)), INK)
    L.gap(0.32)

    L.text([[("Selective PMC as a mechanism test.  ", 22, True, INK),
             ("Correction restricted to the top-P% of dimensions ranked by "
              "|g_{i}|,", 22, False, INK)]], size=22)
    equation(L, [("g_{sel,i} = g_{i}  if  |g_{i}| ≥ |g|_{(P)} ,  0  otherwise ;      "
                  "x′_{sel} = (x + α g_{sel}) / ‖x + α g_{sel}‖", 25, True, INK)],
             note="g_{sel}: truncated gap   ·   |g|_{(P)}: top-P% magnitude "
                  "threshold")
    equation(L, [("E(P)  =  Σ_{i ∈ S(P)} g_{i}²  /  ‖g‖² ,      "
                  "S(P) = { i : |g_{i}| ≥ |g|_{(P)} }", 25, True, INK)],
             note="E(P): captured gap-energy fraction   ·   S(P): top-P% "
                  "index set")
    L.text([[("isolates the mechanism. Since flip risk concentrates in high-|g_{i}| "
              "dimensions, correcting only the highest-energy ones suffices: for "
              "CLIP the top 10% carry 86–92% of gap energy, so P = 5% already "
              "recovers peak recall; ImageBind's diffuse gap (≈72%) requires "
              "full-vector correction.", 22, False, INK)]], size=22)

    # ---------------------------- RIGHT --------------------------------
    R.heading("3", "PMC workflow")
    R.figure("asset_fig2ov.png", 0.86)
    R.gap(0.10)
    R.text([[("Calibration estimates μ_{q}, μ_{x} and g = μ_{q} − μ_{x}; the "
              "offline build shifts and renormalizes database vectors before BQ "
              "index construction; at α = 1 online serving uses q′ = q and "
              "searches the corrected index — zero query-time transform, zero "
              "extra index memory.", 22, False, INK)]], size=22)

    R.heading("4", "Results")
    R.text([[("Five benchmarks (MSCOCO, Flickr30K, Clotho, AudioCaps, LAION-400M — "
              "up to 407 M vectors), three frozen encoders (CLIP-B/32, CLIP-L/14, "
              "ImageBind) and four BQ index types; FAISS IVFRaBitQFastScan at "
              "72–136 B/vec, n_{list} = ⌈√n⌉, n_{probe} ≈ n_{list}/4; R@100 "
              "against exact ground truth, single-thread CPU.",
              21, False, GRAY)]], size=21)
    R.gap(0.12)
    R.text([[("Every configuration, ordered by ‖g‖  ", 24, True, INK),
             ("(Vanilla → PMC)", 20, False, GRAY)]], after=0)
    R.gap(0.12)
    COL_G, R10_W, R100_W, GAPG = (Inches(0.75), Inches(2.35), Inches(2.35),
                                  Inches(0.90))
    GRP_W = R10_W + R100_W
    COL_N = R.w - COL_G - 2 * GRP_W - GAPG
    x0 = R.x + COL_G + COL_N
    x1 = x0 + GRP_W + GAPG
    for gt, gx in (("q→db", x0), ("db→q", x1)):
        cell(slide, [(gt, 21, True, INK)], gx, R.y, GRP_W, Inches(0.40),
             PP_ALIGN.CENTER)
        solid(slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, gx + Inches(0.20), R.y + Inches(0.42),
            GRP_W - Inches(0.40), Pt(1.2)), GRAY)
    R.y += Inches(0.46)
    for t, xx, ww, al in (("‖g‖", R.x, COL_G, PP_ALIGN.LEFT),
                          ("Dataset · Enc.", R.x + COL_G, COL_N, PP_ALIGN.LEFT),
                          ("R@10", x0, R10_W, PP_ALIGN.CENTER),
                          ("R@100", x0 + R10_W, R100_W, PP_ALIGN.CENTER),
                          ("R@10", x1, R10_W, PP_ALIGN.CENTER),
                          ("R@100", x1 + R10_W, R100_W, PP_ALIGN.CENTER)):
        cell(slide, [(t, 20, True, GRAY)], xx, R.y, ww, Inches(0.40), al)
    R.y += Inches(0.42)
    R.rule(Pt(1.5), INK, pad=0.02)
    for (g, name, q10v, q10p, q10d, qv, qp, qd,
         d10v, d10p, d10d, dv, dp, dd) in [
            (".82", "MSCOCO CL-L", ".36", ".48", "+33%", ".55", ".65", "+18%",
             ".26", ".44", "+69%", ".47", ".63", "+34%"),
            (".82", "MSCOCO CLIP", ".40", ".46", "+15%", ".58", ".63", "+9%",
             ".29", ".39", "+34%", ".50", ".60", "+20%"),
            (".77", "Flickr30K CL-L", ".31", ".38", "+23%", ".41", ".48", "+17%",
             ".22", ".38", "+73%", ".33", ".48", "+45%"),
            (".72", "LAION-400M CLIP", ".075", ".086", "+15%", ".108", ".143",
             "+32%", ".035", ".048", "+37%", ".069", ".073", "+6%"),
            (".70", "MSCOCO IB", ".55", ".63", "+15%", ".67", ".75", "+12%",
             ".57", ".64", "+12%", ".71", ".75", "+6%"),
            (".61", "Clotho IB", ".59", ".60", "+2%", ".72", ".73", "+1%",
             ".48", ".54", "+13%", ".62", ".69", "+11%"),
            (".61", "AudioCaps IB", ".39", ".44", "+13%", ".75", ".78", "+4%",
             ".44", ".48", "+9%", ".83", ".83", "+0%")]:
        rh = Inches(0.50)
        cell(slide, [(g, 22, True, CRIMSON)], R.x, R.y, COL_G, rh)
        cell(slide, [(name, 21, False, INK)], R.x + COL_G, R.y, COL_N, rh)
        for vv, pp, dd_, xx, ww in (
                (q10v, q10p, q10d, x0, R10_W),
                (qv, qp, qd, x0 + R10_W, R100_W),
                (d10v, d10p, d10d, x1, R10_W),
                (dv, dp, dd, x1 + R10_W, R100_W)):
            cell(slide, [(f"{vv}→", 17, False, GRAY), (pp, 20, True, INK),
                         (f" {dd_}", 17, True, CRIMSON)],
                 xx, R.y, ww, rh, PP_ALIGN.CENTER)
        R.y += rh
        solid(slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, R.x, R.y, R.w,
                                     Pt(0.8)), HAIR)
    R.gap(0.16)
    R.text([[("Gains scale with ‖g‖: the largest-gap rows improve most, exactly "
              "as Eq. (2) predicts. At 407 M vectors (LAION-400M, 29.3 GB of "
              "codes, 28× compressed) PMC lifts R@100 from .108 to ",
              22, False, INK),
             (".143 ", 22, True, INK), ("(+32%)", 22, True, CRIMSON),
             ("; with exact reranking the lead carries over, .198 → ", 22, False,
              INK),
             (".277 ", 22, True, INK), ("(+40%)", 22, True, CRIMSON),
             (". Exact rescoring cannot rescue Vanilla — it only re-scores what "
              "the corrupted codes already retrieved.", 22, False, INK)]],
           size=22)
    R.gap(0.16)
    R.figure("asset_fig3.png", 1.0)
    R.gap(0.08)
    R.text([[("(a) R@100 vs shift strength α (α = 0: MeanShift, α = 1: full "
              "DB-side PMC): α = 1 is best or near-best, justifying the default. "
              "(b) Selective PMC: on concentrated-gap MSCOCO, top-P = 5% already "
              "reaches peak R@100; diffuse gaps require broader correction. "
              "(c) R@100–QPS Pareto: PMC dominates Vanilla across throughput "
              "levels.", 20, False, GRAY, True)]], size=20, spacing=1.15,
           after=0)

    R.heading("5", "Analysis")
    R.text([[("Every binary index has the same wound  ", 23, True, INK),
             ("(R@100, Vanilla → PMC)", 20, False, GRAY)]], after=0)
    R.gap(0.12)
    T1_W = Emu(int((R.w - Inches(3.2)) / 2))
    for t, xx, ww, al in (("Method", R.x, Inches(3.2), PP_ALIGN.LEFT),
                          ("MSCOCO t→i", R.x + Inches(3.2), T1_W, PP_ALIGN.RIGHT),
                          ("AudioCaps a→t", R.x + Inches(3.2) + T1_W, T1_W,
                           PP_ALIGN.RIGHT)):
        cell(slide, [(t, 21, True, GRAY)], xx, R.y, ww, Inches(0.48), al)
    R.y += Inches(0.50)
    R.rule(Pt(1.5), INK, pad=0.02)
    for m, a, b, c_, d in (("BinaryFlat", ".51", ".57", ".51", ".64"),
                           ("BinaryIVF", ".51", ".57", ".50", ".63"),
                           ("RotatedBinary", ".29", ".58", ".63", ".69"),
                           ("RaBitQ", ".58", ".64", ".67", ".75")):
        rh = Inches(0.50)
        cell(slide, [(m, 22, False, INK)], R.x, R.y, Inches(3.2), rh)
        cell(slide, [(f"{a} → ", 20, False, GRAY), (b, 23, True, INK)],
             R.x + Inches(3.2), R.y, T1_W, rh, PP_ALIGN.RIGHT)
        cell(slide, [(f"{c_} → ", 20, False, GRAY), (d, 23, True, INK)],
             R.x + Inches(3.2) + T1_W, R.y, T1_W, rh, PP_ALIGN.RIGHT)
        R.y += rh
        solid(slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, R.x, R.y, R.w,
                                     Pt(0.8)), HAIR)
    R.gap(0.14)
    R.text([[("Rotation (BBQ-style) spreads the concentrated gap energy across "
              "every bit, so RotatedBinary collapses hardest without correction — "
              "and gains most from it. PMC improves or matches R@100 in all 16 "
              "method × direction configurations: the wound is in the centroid, "
              "so the same one-line fix serves every sign(x − c) method. "
              "Multi-bit quantizers (IVFPQ, OPQ) absorb part of the displacement "
              "yet still gain +1–56%.", 22, False, INK)]], size=22)
    R.gap(0.12)
    R.text([[("Where to correct  ", 24, True, INK),
             ("(MSCOCO · CLIP-B/32 · R@100 t→i)", 20, False, GRAY)]], after=0)
    R.gap(0.14)
    for label, v, col, val in (("Vanilla", 0.578, GRAY, ".578"),
                               ("Query-only", 0.541, GRAY, ".541"),
                               ("Both sides", 0.599, GRAY, ".599"),
                               ("DB-only  (PMC)", 0.637, BLUE, ".637")):
        R.hbar((v - 0.45) / (0.637 - 0.45), col, label, val, h_in=0.50)
    R.gap(0.06)
    R.text([[("Shifting both sides re-displaces the routing pivot, and random, "
              "shuffled or sign-flipped directions of the same norm all fall "
              "below Vanilla — the direction of g is what matters. ",
              22, False, INK),
             ("Correct the centroid where it controls both routing and code "
              "formation.", 22, True, CRIMSON)]], size=22)
    R.gap(0.20)
    R.band([[("Correct the index, not the query:  a one-time, database-side "
              "centroid correction repairs IVF routing and binary quantization "
              "together — zero query-time transform, zero extra index memory, "
              "QPS unchanged at every n_{probe}.", 23, True, WHITE)],
            [("Future work: drift-adaptive α, graph-based ANN, billion-scale "
              "deployment.", 20, False, WHITE)]], pad=0.26)

    print(f"  L: {(L.y - TOP) / 914400.0:5.2f} / {(BOT - TOP) / 914400.0:5.2f} in"
          f"  ({(L.y - TOP) / (BOT - TOP) * 100:4.1f}%)"
          f"  {'OK' if L.y <= BOT else 'OVERFLOW'}")
    print(f"  R: {(R.y - TOP) / 914400.0:5.2f} / {(BOT - TOP) / 914400.0:5.2f} in"
          f"  ({(R.y - TOP) / (BOT - TOP) * 100:4.1f}%)"
          f"  {'OK' if R.y <= BOT else 'OVERFLOW'}")

    # ============================= footer =============================
    fy = PAGE_H - FOOTER_H
    solid(slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, fy, PAGE_W, FOOTER_H), CARD)
    solid(slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, fy, PAGE_W, Pt(3.0)), CRIMSON)
    LOGO_H = Inches(1.30)
    LOGO_W = Emu(int(LOGO_H * ASPECT["asset_cikm.png"]))
    slide.shapes.add_picture(str(HERE / "asset_cikm.png"), M, fy + Inches(0.36),
                             LOGO_W, LOGO_H)
    ft = Col(slide, M + LOGO_W + Inches(0.65), CW - LOGO_W - Inches(0.65),
             fy + Inches(0.50), PAGE_H)
    ft.text([[("The 35th ACM International Conference on Information and Knowledge "
               "Management  ·  November 07–11, 2026  ·  Rome, Italy", 21, True, INK)]],
            after=2)
    ft.text([[("DOI 10.1145/3799682.3840007   ·   CC-BY 4.0   ·   Supported by the "
               "National Research Foundation of Korea (MSIT), No. 00359638   ·   "
               "Code and reproduction package: github.com/sehoon787/PMC",
               18, False, GRAY)]], size=18, after=0)

    out = HERE / "PMC_CIKM2026_poster.pptx"
    prs.save(str(out))
    print(f"saved {out}")


if __name__ == "__main__":
    build()
